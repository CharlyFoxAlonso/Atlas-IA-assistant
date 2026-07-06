"""
core/universal_loader.py
Cargador universal de archivos: PDF, DOCX, PPTX, TXT, MD, Imágenes
Optimizado para manejar archivos grandes sin saturar RAM.
Atlas v3.2
"""
import os
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)  # ✅ CORREGIDO: __name__


def _log(mensaje):
    logger.info(mensaje)


# ============================================
# LECTORES ESPECÍFICOS POR TIPO DE ARCHIVO
# ============================================

def leer_docx(ruta_archivo: str) -> str:
    """Lee un archivo DOCX."""
    try:
        from docx import Document
        doc = Document(ruta_archivo)
        texto = []
        for para in doc.paragraphs:
            if para.text.strip():
                texto.append(para.text)
        return "\n\n".join(texto)
    except Exception as e:
        return f"[Error leyendo DOCX: {str(e)}]"


def leer_pptx(ruta_archivo: str) -> str:
    """Lee un archivo PPTX (PowerPoint)."""
    try:
        from pptx import Presentation
        prs = Presentation(ruta_archivo)
        texto = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():  # ✅ CORREGIDO: sha pe → shape, sin espacios
                    slide_text.append(shape.text)
            if slide_text:
                texto.append(f"\n--- Diapositiva {i} ---\n" + "\n".join(slide_text))
        return "\n".join(texto)
    except Exception as e:
        return f"[Error leyendo PPTX: {str(e)}]"


def leer_imagen(ruta_archivo: str) -> str:
    """Lee una imagen y extrae texto con OCR."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(ruta_archivo)
        texto = pytesseract.image_to_string(img, lang='spa')
        return texto if texto.strip() else "[No se detectó texto en la imagen]"
    except Exception as e:
        return f"[Error leyendo imagen: {str(e)}]"


def leer_texto(ruta_archivo: str) -> str:
    """Lee un archivo de texto plano o Markdown."""
    try:
        with open(ruta_archivo, 'r', encoding='utf-8', errors='ignore') as f:  # ✅ CORREGIDO: sin espacios
            return f.read()
    except Exception as e:
        return f"[Error leyendo texto: {str(e)}]"


def leer_pdf(ruta_archivo: str) -> str:
    """Delega al lector de PDF optimizado."""
    try:
        from core.pdf_reader import leer_pdf as _leer_pdf
        return _leer_pdf(ruta_archivo)
    except Exception as e:
        return f"[Error leyendo PDF: {str(e)}]"


# ============================================
# FUNCIÓN PRINCIPAL: DETECCIÓN AUTOMÁTICA
# ============================================

def leer_archivo(ruta_archivo: str) -> str:
    """
    Detecta el tipo de archivo y lo lee con el método apropiado.
    Soporta: PDF, DOCX, PPTX, TXT, MD, PNG, JPG, JPEG
    """
    if not os.path.exists(ruta_archivo):
        return f"[Error: Archivo no encontrado: {ruta_archivo}]"

    extension = os.path.splitext(ruta_archivo)[1].lower()
    nombre = os.path.basename(ruta_archivo)
    _log(f"Procesando: {nombre} ({extension})")

    if extension == '.pdf':
        return leer_pdf(ruta_archivo)
    elif extension == '.docx':
        return leer_docx(ruta_archivo)
    elif extension == '.pptx':
        return leer_pptx(ruta_archivo)
    elif extension in ['.txt', '.md']:
        return leer_texto(ruta_archivo)
    elif extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
        return leer_imagen(ruta_archivo)
    else:
        return f"[Tipo de archivo no soportado: {extension}]"


def leer_archivo_con_info(ruta_archivo: str) -> dict:
    """
    Lee un archivo y devuelve metadata + contenido.
    CLAVE: Usa 'ok' en lugar de 'exito' para compatibilidad con indexer.py
    """
    if not os.path.exists(ruta_archivo):
        return {
            "ok": False,
            "error": f"Archivo no encontrado: {ruta_archivo}"
        }
    try:
        contenido = leer_archivo(ruta_archivo)

        # Verificar que el contenido no sea un error
        if contenido.startswith("[Error") or contenido.startswith("[Tipo de archivo no soportado"):
            return {
                "ok": False,
                "error": contenido
            }

        stat = os.stat(ruta_archivo)
        return {
            "ok": True,  # ✅ CAMBIO CLAVE: 'ok' en lugar de 'exito'
            "ruta": ruta_archivo,
            "nombre": os.path.basename(ruta_archivo),
            "contenido": contenido,
            "tamano_bytes": stat.st_size,
            "tamano_kb": round(stat.st_size / 1024, 2)
        }
    except Exception as e:
        return {
            "ok": False,
            "error": str(e)
        }


def mostrar_dependencias():
    """Muestra qué dependencias están instaladas."""
    deps = {
        "pypdf": "Lectura de PDFs",
        "pdf2image": "Conversión PDF a imagen (OCR)",
        "pytesseract": "OCR (reconocimiento de texto)",
        "PIL": "Procesamiento de imágenes",
        "python-docx": "Lectura de DOCX",
        "python-pptx": "Lectura de PPTX"
    }

    print("\n" + "=" * 60)
    print("DEPENDENCIAS DE UNIVERSAL LOADER")
    print("=" * 60)
    for modulo, descripcion in deps.items():
        try:
            __import__(modulo)
            print(f"✅ {modulo:20s} - {descripcion}")
        except ImportError:
            print(f"❌ {modulo:20s} - {descripcion} [FALTA]")
    print("=" * 60 + "\n")


if __name__ == "__main__":  # ✅ CORREGIDO: __name__ == "__main__"
    import sys
    if len(sys.argv) > 1:
        archivo = sys.argv[1]
        print(f"\nProbando lectura de: {archivo}\n")
        resultado = leer_archivo(archivo)
        print(f"\nContenido extraído ({len(resultado)} caracteres):")
        print("-" * 60)
        print(resultado[:2000])
        if len(resultado) > 2000:
            print(f"\n... ({len(resultado) - 2000} caracteres más)")
    else:
        mostrar_dependencias()